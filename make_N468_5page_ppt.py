from __future__ import annotations

import math
from pathlib import Path
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    print("未检测到 python-pptx。请先安装：pip install python-pptx")
    sys.exit(1)


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True


def add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int = 16, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(size)
            run.font.bold = bold


def add_labeled_box(
    slide,
    title: str,
    subtitle: str,
    left: float,
    top: float,
    width: float,
    height: float,
    title_size: int = 14,
    subtitle_size: int = 11,
) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(242, 242, 242)
    rect.line.color.rgb = RGBColor(120, 120, 120)
    add_text(slide, title, left + 0.12, top + 0.08, width - 0.24, 0.28, size=title_size, bold=True)
    add_text(slide, subtitle, left + 0.12, top + 0.38, width - 0.24, height - 0.46, size=subtitle_size)


def draw_top_view_panel(slide, left: float, top: float, width: float, height: float, n_cyl: int) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(250, 250, 250)
    panel.line.color.rgb = RGBColor(150, 150, 150)
    add_text(slide, f"N={n_cyl}", left + 0.08, top + 0.05, 0.8, 0.25, size=14, bold=True)

    cx, cy = left + width * 0.5, top + height * 0.52
    r = min(width, height) * 0.26
    dot = 0.09
    for i in range(n_cyl):
        theta = 2 * math.pi * i / n_cyl
        x = cx + r * math.cos(theta) - dot / 2
        y = cy + r * math.sin(theta) - dot / 2
        c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(dot), Inches(dot))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(34, 102, 170)
        c.line.color.rgb = RGBColor(34, 102, 170)

    probe_map = {
        "P0": (0.0, 0.0),
        "P1": (0.12, 0.0),
        "P2": (-0.12, 0.0),
        "P3": (0.0, 0.12),
        "P4": (0.0, -0.12),
    }
    for name, (px, py) in probe_map.items():
        x = cx + px - 0.03
        y = cy + py - 0.03
        p = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.06), Inches(0.06))
        p.fill.solid()
        p.fill.fore_color.rgb = RGBColor(200, 70, 70)
        p.line.color.rgb = RGBColor(200, 70, 70)
        add_text(slide, name, x + 0.03, y - 0.08, 0.28, 0.16, size=9)

    add_text(slide, "front (-2,0)", left + 0.08, cy - 0.10, 0.8, 0.18, size=9)
    add_text(slide, "rear (2,0)", left + width - 0.85, cy - 0.10, 0.7, 0.18, size=9)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(cx - 0.65), Inches(top + height - 0.34), Inches(1.3), Inches(0.16))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(90, 90, 90)
    arrow.line.color.rgb = RGBColor(90, 90, 90)
    add_text(slide, "波向: front → center → rear", left + 0.35, top + height - 0.18, width - 0.7, 0.15, size=8)


def add_image_keep_ratio(slide, image_path: Path, left: float, top: float, width: float, height: float) -> bool:
    if not image_path.exists():
        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        add_text(slide, f"缺少图片:\n{image_path.as_posix()}", left + 0.1, top + 0.1, width - 0.2, height - 0.2, size=14)
        return False

    pic = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    max_h = Inches(height)
    if pic.height > max_h:
        scale = max_h / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
    pic.left = Inches(left) + int((Inches(width) - pic.width) / 2)
    pic.top = Inches(top) + int((Inches(height) - pic.height) / 2)
    return True


def main() -> int:
    base_dir = Path(__file__).resolve().parent
    figures = base_dir / "figures"
    output_path = base_dir / "outputs" / "N468_5page_results_summary.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Page 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "研究对象与测点定义")
    add_text(
        s1,
        "采用统一几何参数、统一测点定义和统一扫频范围，对 N=4/6/8 单环圆柱阵列进行横向比较。",
        0.6,
        6.6,
        12.0,
        0.6,
        size=16,
    )
    box_w, box_h, y0 = 4.15, 4.25, 1.1
    for i, n_cyl in enumerate([4, 6, 8]):
        x = 0.35 + i * 4.33
        draw_top_view_panel(s1, x, y0, box_w, box_h, n_cyl)
    add_text(
        s1,
        "三组模型保持圆柱半径、水深、环半径和吃水一致，通过改变环向圆柱数量比较不同环向离散度和多体散射路径下的中心响应差异。",
        0.55,
        6.15,
        12.2,
        0.9,
        size=15,
    )

    # Page 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "数据流程与指标定义")
    add_labeled_box(s2, "Capytaine-BEM frequency scan", "run_ring_array.py", 0.55, 1.2, 3.0, 1.15)
    add_labeled_box(s2, "Full scan CSV", "analyze_scan_results.py / plot_scan_diagnostics.py", 3.75, 1.2, 2.7, 1.15)
    add_labeled_box(s2, "Single-model diagnostics", "compare_N4_N6_N8.py", 6.65, 1.2, 2.95, 1.15)
    add_labeled_box(s2, "N=4/6/8 comparison figures and tables", "make_N468_5page_ppt.py", 9.8, 1.2, 3.0, 1.15)
    for ax in [3.58, 6.48, 9.62]:
        ar = s2.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(ax), Inches(1.55), Inches(0.15), Inches(0.35))
        ar.fill.solid()
        ar.fill.fore_color.rgb = RGBColor(90, 90, 90)
        ar.line.color.rgb = RGBColor(90, 90, 90)
    defs = (
        "指标定义:\n"
        "• 所有响应指标均基于频域总自由面复振幅模值 total_abs = |η_incident + η_diffracted|\n"
        "• center_mean_abs：P0–P4 平均响应\n"
        "• center_max_abs：P0–P4 局部最大响应\n"
        "• center_max_to_mean_ratio：中心局部化指标\n"
        "• rear_abs：背浪侧测点响应\n"
        "• S_rear_front = rear_abs / front_abs，仅为 transmission-like indicator，不是 Kt。"
    )
    add_text(s2, defs, 0.75, 2.65, 12.0, 3.9, size=15)

    # Page 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "N=4/6/8 中心响应总对比")
    add_image_keep_ratio(s3, figures / "N468_center_mean_comparison.png", 0.45, 0.95, 6.35, 3.45)
    add_image_keep_ratio(s3, figures / "N468_center_max_comparison.png", 6.6, 0.95, 6.35, 3.45)
    add_text(s3, "N | center_mean peak | center_max peak", 0.95, 4.55, 6.7, 0.35, size=13, bold=True)
    add_text(s3, "4 | 1.139 @ 0.96 s | 1.181 @ 0.90 s*", 0.95, 4.9, 6.7, 0.3, size=13)
    add_text(s3, "6 | 1.238 @ 0.97 s | 1.285 @ 0.90 s*", 0.95, 5.18, 6.7, 0.3, size=13)
    add_text(s3, "8 | 1.355 @ 1.00 s | 1.419 @ 0.94 s", 0.95, 5.46, 6.7, 0.3, size=13)
    add_text(s3, "注：* 表示 T=0.90 s 边界峰。", 0.95, 5.78, 4.8, 0.28, size=12)
    add_text(
        s3,
        "当前 0.90–2.00 s 统一后处理口径下，中心平均响应和局部最大响应均表现为 N=8 > N=6 > N=4；其中 N=4/N=6 的 center_max 峰值属于 T=0.90 s 边界峰。",
        0.55,
        6.15,
        12.2,
        0.85,
        size=14,
        bold=True,
    )

    # Page 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "中心局部化与 P0–P4 空间分化")
    add_image_keep_ratio(s4, figures / "N468_center_max_to_mean_ratio_comparison.png", 0.55, 1.0, 7.55, 4.95)
    add_image_keep_ratio(s4, figures / "N468_center_probe_peak_period_comparison.png", 8.25, 1.0, 4.45, 2.35)
    add_image_keep_ratio(s4, figures / "N468_center_probe_peak_comparison.png", 8.25, 3.6, 4.45, 2.35)
    add_text(s4, "N=8 在短周期端表现出更高的局部化指标；P0–P4 峰值周期并不同步，表明中心增强具有空间分化特征，而非均匀同步增强。", 0.7, 6.15, 12.2, 0.9, size=15, bold=True)

    # Page 5
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s5, "前后场指标、边界峰风险与下一步")
    add_image_keep_ratio(s5, figures / "N468_front_rear_S_comparison.png", 0.6, 1.2, 7.0, 4.8)
    add_labeled_box(s5, "指标边界", "S_rear_front = rear_abs / front_abs，仅为 transmission-like indicator，不是 Kt。", 7.85, 1.3, 4.95, 1.35, title_size=15, subtitle_size=13)
    add_labeled_box(s5, "边界峰风险", "多个关键峰值位于 T=0.90 s，短周期端物理解释需保守。", 7.85, 2.95, 4.95, 1.35, title_size=15, subtitle_size=13)
    add_labeled_box(s5, "下一步", "补充 T=0.85–1.05 s，ΔT=0.005 s 加密扫描。", 7.85, 4.6, 4.95, 1.35, title_size=15, subtitle_size=13)

    prs.save(output_path)
    print(str(output_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
