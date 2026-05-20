from __future__ import annotations

from pathlib import Path
import sys

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("未检测到 python-pptx。请先安装：pip install python-pptx")
    sys.exit(1)


TITLE = "柱数变化对中心响应与局部化特征的影响"
BOTTOM_SUMMARY = (
    "在当前统一后处理口径下，N=8 阵列表现出更高的中心响应和局部化特征；"
    "但短周期边界峰值仍需保守解释。"
)

FIG_ITEMS = [
    (
        "figures/N468_center_mean_comparison.png",
        "中心平均响应：N=8 > N=6 > N=4",
    ),
    (
        "figures/N468_center_max_comparison.png",
        "中心局部最大响应：N=8 最高，但 N=4/N=6 峰值位于 0.90 s 边界",
    ),
    (
        "figures/N468_center_max_to_mean_ratio_comparison.png",
        "局部化程度：N=8 的 max/mean ratio 最高，中心热点更突出",
    ),
    (
        "figures/N468_center_probe_peak_period_comparison.png",
        "空间分化：P1 峰值最早，P2 峰值最晚，P3/P4 对称",
    ),
]


def add_centered_textbox(slide, text, left, top, width, height, font_size=14, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    return box


def main() -> int:
    base_dir = Path(__file__).resolve().parent
    output_path = base_dir / "outputs" / "N468_summary_comparison_slide.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    for fig_path, _ in FIG_ITEMS:
        full_path = base_dir / fig_path
        if not full_path.exists():
            print(f"缺少输入图片：{full_path}")
            return 1

    prs = Presentation()
    # 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_centered_textbox(
        slide,
        TITLE,
        Inches(0.3),
        Inches(0.1),
        Inches(12.7),
        Inches(0.5),
        font_size=24,
        bold=True,
    )

    margin_x = 0.4
    top_start = 0.75
    col_gap = 0.35
    row_gap = 0.18
    caption_h = 0.42
    summary_h = 0.7
    summary_margin_bottom = 0.15

    usable_w = 13.333 - 2 * margin_x - col_gap
    panel_w = usable_w / 2

    summary_top = 7.5 - summary_h - summary_margin_bottom
    panels_total_h = summary_top - top_start
    panel_h = (panels_total_h - row_gap) / 2
    image_h = panel_h - caption_h

    positions = [
        (margin_x, top_start),
        (margin_x + panel_w + col_gap, top_start),
        (margin_x, top_start + panel_h + row_gap),
        (margin_x + panel_w + col_gap, top_start + panel_h + row_gap),
    ]

    for (fig_rel, caption), (left, top) in zip(FIG_ITEMS, positions):
        img_path = base_dir / fig_rel
        pic = slide.shapes.add_picture(
            str(img_path),
            Inches(left),
            Inches(top),
            width=Inches(panel_w),
        )
        if pic.height > Inches(image_h):
            scale = Inches(image_h) / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
        pic.left = Inches(left) + int((Inches(panel_w) - pic.width) / 2)
        pic.top = Inches(top) + int((Inches(image_h) - pic.height) / 2)
        add_centered_textbox(
            slide,
            caption,
            Inches(left),
            Inches(top + image_h),
            Inches(panel_w),
            Inches(caption_h),
            font_size=12,
        )

    add_centered_textbox(
        slide,
        BOTTOM_SUMMARY,
        Inches(0.4),
        Inches(summary_top),
        Inches(12.533),
        Inches(summary_h),
        font_size=13,
        bold=True,
    )

    prs.save(output_path)
    print(str(output_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
